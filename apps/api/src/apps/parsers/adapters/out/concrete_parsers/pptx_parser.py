"""
PPTX парсер для извлечения текста и изображений из .pptx файлов.

Использует библиотеку python-pptx для извлечения текста, таблиц и изображений из презентаций.
"""
import logging
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ....domain.ports import DocumentParser
from ....domain.models import ParsedDocument, DocumentImage

logger = logging.getLogger(__name__)


class PptxDocumentParser(DocumentParser):
    """
    Парсер для PPTX файлов.
    
    Извлекает:
    - Текст из слайдов (заголовки, текстовые блоки, таблицы)
    - Структуру презентации (заголовки слайдов)
    - Изображения из презентации
    """

    def __init__(
        self,
        min_width: int = 50,
        min_height: int = 50,
        min_file_size: int = 3 * 1024,
    ):
        """
        Args:
            min_width: Минимальная ширина изображения в пикселях
            min_height: Минимальная высота изображения в пикселях
            min_file_size: Минимальный размер файла изображения в байтах
        """
        self.min_width = min_width
        self.min_height = min_height
        self.min_file_size = min_file_size

    def parse(
        self,
        file_bytes: bytes,
        file_name: str,
        process_images: bool = False,
    ) -> ParsedDocument:
        """
        Парсит PPTX документ.

        Args:
            file_bytes: Байты PPTX файла
            file_name: Имя файла (для логирования)
            process_images: Извлекать ли изображения

        Returns:
            ParsedDocument с результатом парсинга

        Raises:
            Exception: Если не удается обработать PPTX файл
        """
        try:
            logger.info(f"📊 Парсинг PPTX: {file_name}")
            
            # Открываем презентацию из байтов
            prs_stream = BytesIO(file_bytes)
            prs = Presentation(prs_stream)
            
            logger.info(f"PPTX открыт. Количество слайдов: {len(prs.slides)}")
            
            # Извлекаем текст из всех слайдов
            text_parts = []
            images: list[DocumentImage] = []
            contents = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                logger.debug(f"Обработка слайда {slide_num}/{len(prs.slides)}")
                
                # Извлекаем текст из слайда
                slide_text = self.extract_text_from_slide(slide, slide_num)
                text_parts.append(slide_text)
                
                # Добавляем заголовок слайда в оглавление
                slide_title = self.get_slide_title(slide)
                if slide_title:
                    contents.append({
                        "level": 1,
                        "title": slide_title,
                        "page": slide_num,
                    })
                
                # Извлекаем изображения если нужно
                if process_images:
                    slide_images = self.extract_images_from_slide(slide, slide_num)
                    images.extend(slide_images)

                
            
            final_text = "\n\n".join(text_parts)
            
            logger.info(
                f"✅ PPTX парсинг завершен. Слайдов: {len(prs.slides)}, "
                f"Изображений: {len(images)}"
            )
            
            return ParsedDocument(
                text=final_text,
                images=images,
                contents=contents,
                is_book=False,  # PPTX обычно не являются книгами
            )
            
        except Exception as e:
            logger.error(f"❌ Ошибка при парсинге PPTX файла {file_name}: {str(e)}")
            raise Exception(f"Ошибка при парсинге PPTX файла: {str(e)}")
    
    def get_slide_title(self, slide) -> str:
        """
        Получает заголовок слайда.
        
        Args:
            slide: Объект Slide из python-pptx
            
        Returns:
            Заголовок слайда или пустая строка
        """
        # Проверяем shapes.title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
        
        # Ищем первый текстовый блок как заголовок
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    return text
        
        return ""
    
    def extract_text_from_slide(self, slide, slide_num: int) -> str:
        """
        Извлекает текст из слайда.
        
        Обрабатывает текстовые блоки, таблицы и другие элементы.
        
        Args:
            slide: Объект Slide из python-pptx
            slide_num: Номер слайда
            
        Returns:
            Текст слайда
        """
        text_parts = []
        
        for shape in slide.shapes:
            # Текстовые блоки
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
            
            # Таблицы
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_text = self.format_table(shape.table)
                if table_text:
                    text_parts.append(table_text)
            
            # Группы фигур (рекурсивно)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = self.extract_text_from_group_shape(shape)
                if group_text:
                    text_parts.append(group_text)
        
        return "\n\n".join(text_parts)
    
    def extract_text_from_group_shape(self, group_shape) -> str:
        """
        Извлекает текст из группы фигур.
        
        Args:
            group_shape: Группа фигур
            
        Returns:
            Текст из группы фигур
        """
        text_parts = []
        
        for shape in group_shape.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
            
            # Рекурсивно обрабатываем вложенные группы
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                nested_text = self.extract_text_from_group_shape(shape)
                if nested_text:
                    text_parts.append(nested_text)
        
        return "\n".join(text_parts)
    
    def format_table(self, table) -> str:
        """
        Форматирует таблицу в текстовый вид.
        
        Args:
            table: Объект Table из python-pptx
            
        Returns:
            Отформатированная таблица
        """
        if not table.rows:
            return ""
        
        # Собираем все строки таблицы
        table_rows = []
        for row in table.rows:
            cell_texts = [cell.text.strip() for cell in row.cells]
            table_rows.append(cell_texts)
        
        if not table_rows:
            return ""
        
        # Вычисляем ширину колонок
        num_cols = max(len(row) for row in table_rows)
        col_widths = [0] * num_cols
        
        for row in table_rows:
            for i, cell in enumerate(row):
                if i < num_cols:
                    col_widths[i] = max(col_widths[i], len(cell))
        
        # Форматируем таблицу
        formatted_rows = []
        for i, row in enumerate(table_rows):
            # Дополняем строку пустыми ячейками если нужно
            while len(row) < num_cols:
                row.append("")
            
            # Форматируем ячейки с выравниванием
            formatted_cells = [
                cell.ljust(col_widths[j]) for j, cell in enumerate(row)
            ]
            formatted_rows.append("| " + " | ".join(formatted_cells) + " |")
            
            # Добавляем разделитель после заголовка (первой строки)
            if i == 0:
                separator = "|" + "|".join(["-" * (w + 2) for w in col_widths]) + "|"
                formatted_rows.append(separator)
        
        return "\n".join(formatted_rows)
    
    def extract_images_from_slide(self, slide, slide_num: int) -> list[DocumentImage]:
        """
        Извлекает изображения из слайда.
        
        Args:
            slide: Объект Slide из python-pptx
            slide_num: Номер слайда
            
        Returns:
            Список изображений
        """
        images = []
        image_index = 0
        
        for shape in slide.shapes:
            # Изображения
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Проверяем размер файла
                    if len(image_bytes) < self.min_file_size:
                        logger.debug(
                            f"Пропуск изображения на слайде {slide_num}: "
                            f"размер {len(image_bytes)} < {self.min_file_size} байт"
                        )
                        continue
                    
                    # Определяем расширение
                    ext = image.ext or "png"
                    if ext.startswith('.'):
                        ext = ext[1:]
                    
                    # Получаем размеры (в пикселях, приблизительно)
                    # python-pptx возвращает размеры в EMU (English Metric Units)
                    # 1 дюйм = 914400 EMU, предполагаем 96 DPI
                    width = int(shape.width / 914400 * 96) if hasattr(shape, 'width') else 0
                    height = int(shape.height / 914400 * 96) if hasattr(shape, 'height') else 0
                    
                    # Проверяем минимальные размеры
                    if width < self.min_width or height < self.min_height:
                        logger.debug(
                            f"Пропуск изображения на слайде {slide_num}: "
                            f"размер {width}x{height} меньше минимального "
                            f"{self.min_width}x{self.min_height}"
                        )
                        continue
                    
                    doc_image = DocumentImage(
                        bytes=image_bytes,
                        ext=ext,
                        width=width,
                        height=height,
                        page=slide_num,
                        index=image_index,
                        marker=f"[IMAGE_{slide_num}_{image_index}]",
                        file_name=f"slide_{slide_num}_image_{image_index}.{ext}",
                    )
                    
                    images.append(doc_image)
                    image_index += 1
                    logger.debug(
                        f"Извлечено изображение со слайда {slide_num}: "
                        f"{width}x{height}, {len(image_bytes)} байт"
                    )
                    
                except Exception as e:
                    logger.warning(
                        f"Не удалось извлечь изображение со слайда {slide_num}: {str(e)}"
                    )
            
            # Обрабатываем группы фигур
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_images = self.extract_images_from_group_shape(
                    shape, slide_num, image_index
                )
                images.extend(group_images)
                image_index += len(group_images)
        
        return images
    
    def extract_images_from_group_shape(
        self, group_shape, slide_num: int, start_index: int
    ) -> list[DocumentImage]:
        """
        Извлекает изображения из группы фигур.
        
        Args:
            group_shape: Группа фигур
            slide_num: Номер слайда
            start_index: Начальный индекс для изображений
            
        Returns:
            Список изображений
        """
        images = []
        image_index = start_index
        
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    if len(image_bytes) < self.min_file_size:
                        continue
                    
                    ext = image.ext or "png"
                    if ext.startswith('.'):
                        ext = ext[1:]
                    
                    width = int(shape.width / 914400 * 96) if hasattr(shape, 'width') else 0
                    height = int(shape.height / 914400 * 96) if hasattr(shape, 'height') else 0
                    
                    if width < self.min_width or height < self.min_height:
                        continue
                    
                    doc_image = DocumentImage(
                        bytes=image_bytes,
                        ext=ext,
                        width=width,
                        height=height,
                        page=slide_num,
                        index=image_index,
                        marker=f"[IMAGE_{slide_num}_{image_index}]",
                        file_name=f"slide_{slide_num}_image_{image_index}.{ext}",
                    )
                    
                    images.append(doc_image)
                    image_index += 1
                    
                except Exception as e:
                    logger.warning(
                        f"Не удалось извлечь изображение из группы на слайде {slide_num}: {str(e)}"
                    )
            
            # Рекурсивно обрабатываем вложенные группы
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                nested_images = self.extract_images_from_group_shape(
                    shape, slide_num, image_index
                )
                images.extend(nested_images)
                image_index += len(nested_images)
        
        return images
